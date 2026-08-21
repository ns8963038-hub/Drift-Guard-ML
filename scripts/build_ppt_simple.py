"""Generate the presentation-friendly version of the deck.

The first deck reads well and presents badly: 112 bullets averaging 16 words,
the longest 38. Bullets that long are complete sentences, so a presenter has no
choice but to read them out, and the audience reads ahead instead of listening.

This version inverts that. The slide carries a short cue — a few words the
presenter glances at — and the sentence they actually say lives in the speaker
notes, where the audience never sees it. Same 17-section structure, same
verified figures, same black-on-white treatment.

Layout helpers are shared with build_ppt.py so the two decks cannot drift apart
stylistically.

    python scripts/build_ppt_simple.py    # -> dist/DriftGuard_Presentation_Simple.pptx
"""

from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

from build_ppt import (  # noqa: E402
    BODY_W,
    GREY,
    MARGIN,
    RULE,
    SLIDE_H,
    SLIDE_W,
    _blank,
    _text,
    heading,
    table,
)

OUT = (
    Path(__file__).resolve().parent.parent
    / "dist"
    / "DriftGuard_Presentation_Simple.pptx"
)


def cues(slide, items, top=Inches(2.0), size=22):
    """Short bullets, generously spaced. No sentence should reach two lines."""
    box = slide.shapes.add_textbox(MARGIN, top, BODY_W, SLIDE_H - top - Inches(0.5))
    frame = box.text_frame
    frame.word_wrap = True

    for i, item in enumerate(items):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        para.space_after = Pt(16)
        para.line_spacing = 1.1
        run = para.add_run()
        run.text = f"•   {item}"
        run.font.size = Pt(size)
        run.font.name = "Calibri"


def notes(slide, text):
    """What the presenter says. Never shown to the audience."""
    slide.notes_slide.notes_text_frame.text = text.strip()


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── Title ────────────────────────────────────────────────────────
    s = _blank(prs)
    _, run = _text(
        s,
        MARGIN,
        Inches(2.4),
        BODY_W,
        Inches(1.2),
        44,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    run.text = "DriftGuard"
    _, run = _text(
        s, MARGIN, Inches(3.35), BODY_W, Inches(1.0), 22, align=PP_ALIGN.CENTER
    )
    run.text = "Watching Machine Learning Models\nAfter They Go Live"
    line = s.shapes.add_shape(1, Inches(5.4), Inches(4.6), Inches(2.5), Pt(1.2))
    line.fill.solid()
    line.fill.fore_color.rgb = RULE
    line.line.fill.background()
    line.shadow.inherit = False
    _, run = _text(
        s,
        MARGIN,
        Inches(5.0),
        BODY_W,
        Inches(1.2),
        15,
        color=GREY,
        align=PP_ALIGN.CENTER,
    )
    run.text = "Final Year Project\n\n[ Student names  ·  USN  ·  Guide  ·  Department  ·  College  ·  Year ]"
    notes(
        s,
        """
Good morning. Our project is called DriftGuard.

A machine learning model is trained once, on data from the past. Then it is put
into use, and the world keeps changing around it. The model does not change.

Our project watches for that, and warns you before it costs you anything.
""",
    )

    # ── 1 ────────────────────────────────────────────────────────────
    s = _blank(prs)
    heading(s, "1", "Introduction")
    cues(
        s,
        [
            "A model is trained once, on old data",
            "The world keeps changing. The model does not",
            "It never crashes — it just becomes wrong",
            "DriftGuard watches for this and warns you",
            "Three types of user, with different powers",
        ],
    )
    notes(
        s,
        """
A machine learning model learns patterns from old data. Once trained, it is
fixed. It cannot learn anything new by itself.

But the world it is used in keeps changing. Customers change, prices change,
behaviour changes. The model does not notice any of this.

Here is the important part. When a model goes wrong, it does not crash. There
is no error message. It keeps giving answers, confidently, and those answers
are simply wrong. Nobody finds out until the business loses money.

DriftGuard is a website that watches a live model, checks every new batch of
data, and raises a warning when something changes.

It also supports three types of user, because in a real company different
people have different responsibilities.
""",
    )

    # ── 2 ────────────────────────────────────────────────────────────
    s = _blank(prs)
    heading(s, "2", "Motivation")
    cues(
        s,
        [
            "Normal software crashes. Models fail silently",
            "A churn model stops finding customers who leave",
            "Checking 19 columns by hand does not scale",
            "The true answers arrive weeks later",
            "We need a warning from the data itself",
        ],
    )
    notes(
        s,
        """
When normal software fails, you know immediately. It crashes, or throws an
error. A machine learning model gives you nothing.

Take a real example. A telecom company uses a model to find customers about to
cancel, so they can call them and offer a discount. If that model quietly stops
working, the company keeps spending its retention budget on the wrong people
and never knows why results dropped.

Somebody could check the data by hand, but that means checking nineteen columns,
for every model, for every batch, every day. That is not realistic.

And you cannot simply wait for accuracy to tell you, because the true answer
comes much later. Whether a customer cancels is only known weeks or months
afterwards.

So we need a system that spots the problem in the incoming data itself, before
the true answers arrive.
""",
    )

    # ── 3 ────────────────────────────────────────────────────────────
    s = _blank(prs)
    heading(s, "3", "Problem Statement")
    cues(
        s,
        [
            "Models get worse as data changes over time",
            "This is called data drift",
            "Nothing warns you when it happens",
            "True answers come too late to rely on",
            "Different people need different access",
        ],
    )
    notes(
        s,
        """
Our problem statement is this.

Machine learning models deployed in production get worse over time, because the
incoming data slowly stops looking like the data they were trained on. That is
called data drift.

Today, nothing warns you when this happens. The failure is silent, and it is
usually discovered only after real business loss.

Three things make the problem harder. First, the true answers arrive late, so
accuracy cannot be the main signal. Second, statistical tests alone are not
trustworthy, because their result depends heavily on how many rows you have.
Third, in a real company several people are involved, and they should not all
have the same powers.

So we need a system that detects drift from the input data alone, explains it
simply, and controls who is allowed to change what.
""",
    )

    # ── 4 ────────────────────────────────────────────────────────────
    s = _blank(prs)
    heading(s, "4", "Objectives")
    cues(
        s,
        [
            "Detect drift using four different measures",
            "Check the quality of incoming data",
            "Track accuracy, and handle missing answers",
            "Give one health score out of 100",
            "Raise clear alerts in plain English",
            "Advise retraining — never do it automatically",
            "Enforce three user roles on the server",
        ],
        size=20,
    )
    notes(
        s,
        """
These are our objectives.

First, detect drift using four different measures, so that we are not depending
on any single test.

Second, check the quality of the incoming data — missing values, duplicate
rows, outliers, and categories we have never seen before.

Third, track how accurate the model is over time, and handle correctly the very
common case where the true answers are not available yet.

Fourth, combine everything into one health score out of 100, and show how that
score was built up.

Fifth, raise alerts written in plain English, not in statistics.

Sixth, recommend retraining, but never perform it automatically. That decision
belongs to a person.

And seventh, implement three user roles, enforced by the server itself.
""",
    )

    # ── 5 ────────────────────────────────────────────────────────────
    s = _blank(prs)
    heading(s, "5", "Literature Survey")
    table(
        s,
        ["Work", "What it gave us"],
        [
            [
                "Sculley et al., 2015",
                "Showed monitoring is the biggest hidden cost of deployed models",
            ],
            [
                "Gama et al., 2014",
                "The standard survey of drift; gave us the terminology",
            ],
            [
                "Lu et al., 2019",
                "Compared drift detection methods; suggested using several together",
            ],
            ["Massey, 1951", "The Kolmogorov-Smirnov test we use on number columns"],
            ["Lin, 1991", "Jensen-Shannon Divergence, one of our two size measures"],
            ["Siddiqi, 2006", "PSI and its 0.10 / 0.25 bands, from credit scoring"],
            [
                "Breck et al., 2017",
                "Argued models must keep being tested after release",
            ],
            [
                "Polyzotis et al., 2017",
                "Argued input data must be validated, not trusted",
            ],
        ],
        col_widths=[2.4, 7.6],
        size=14,
        top=Inches(1.95),
    )
    notes(
        s,
        """
We studied eight main works.

Sculley and colleagues, in 2015, showed that monitoring is the largest hidden
cost of running machine learning in production. That paper is the reason this
project exists.

Gama and colleagues, and later Lu and colleagues, surveyed how drift is
detected. Lu's review is where we got the idea of using several measures
together rather than trusting one.

Massey gave us the Kolmogorov-Smirnov test for number columns. Lin gave us
Jensen-Shannon Divergence. Siddiqi's book on credit scoring is where the
Population Stability Index comes from, along with the 0.10 and 0.25 bands we
use.

Breck and Polyzotis both argued that after a model is released, its data and
its behaviour must keep being checked. That is exactly what we built.
""",
    )

    # ── 6 ────────────────────────────────────────────────────────────
    s = _blank(prs)
    heading(s, "6", "Research Gap")
    cues(
        s,
        [
            "Research studies algorithms, not usable systems",
            "Commercial tools are costly and cloud-only",
            "Most rely on statistical tests alone",
            "Missing answers are wrongly counted as zero",
            "Access control is rarely considered",
        ],
    )
    notes(
        s,
        """
So what is missing?

Most research studies the drift algorithms themselves, in isolation. Very
little of it becomes a system that a team can actually sit down and use.

There are commercial products, but they are expensive, closed, and require the
cloud. That puts them out of reach for a college project or a small company.

Three specific gaps. Most tools decide drift using a statistical test alone,
even though the result of that test depends mostly on how many rows you have,
not on how important the change is.

Many treat a batch with no answers as scoring zero accuracy, which raises false
alarms about models that are perfectly healthy.

And almost none of them consider access control, even though in practice the
person reading the alerts should not be the person who can change the alert
settings.
""",
    )

    # ── 7 ────────────────────────────────────────────────────────────
    s = _blank(prs)
    heading(s, "7", "Proposed Solution")
    cues(
        s,
        [
            "A Django website you host yourself",
            "Every new batch is compared to a fixed baseline",
            "Size of change decides; the test only confirms",
            "No answers? Record unknown, not zero",
            "Settings are saved with each run",
            "Three roles, enforced by the server",
        ],
        size=20,
    )
    notes(
        s,
        """
Our solution is a Django website that a college or a company can host on its own
machine. No cloud, no subscription.

When we register a model, we give it a baseline — a fixed reference dataset,
usually the data it was trained on. Every new batch that arrives is compared
against that baseline.

Now the key idea. We decide the drift level by how big the change is, using PSI
and Jensen-Shannon Divergence. The statistical test only confirms it. We do not
let the test decide, because on a very large batch even a meaningless change
looks significant, and on a small batch even a serious change does not.

If a batch arrives without answers, we record accuracy as unknown, not zero, and
the health score works from the parts it can still measure.

And the threshold settings in force are saved with every run. So if somebody
changes a setting next month, old results still show what they were actually
judged against.
""",
    )

    # ── 8 ────────────────────────────────────────────────────────────
    s = _blank(prs)
    heading(s, "8", "Hardware and Software")
    table(
        s,
        ["", "What we used"],
        [
            ["Computer", "Any laptop — 4 GB RAM, 2 GB free disk"],
            ["Operating system", "Windows, macOS or Linux"],
            ["Language", "Python 3.11"],
            ["Web framework", "Django 5.0.9"],
            ["Machine learning", "scikit-learn, joblib"],
            ["Data and statistics", "pandas, NumPy, SciPy"],
            ["Database", "SQLite"],
            ["Charts", "Chart.js (stored locally, no internet needed)"],
            ["Background jobs", "APScheduler"],
            ["Testing", "pytest"],
        ],
        col_widths=[3.0, 7.0],
        size=15,
        top=Inches(1.95),
    )
    notes(
        s,
        """
The requirements are deliberately small. It runs on any ordinary laptop, with
4 GB of RAM and 2 GB of free disk space. It works on Windows, macOS or Linux.

We used Python 3.11 — specifically 3.11, because Django 5 does not support
Python 3.12 or newer.

Django is the web framework. scikit-learn trains and runs the models. pandas,
NumPy and SciPy do the data handling and the statistics.

The database is SQLite, which is a single file, so there is nothing to install.

The charts use Chart.js, and we store it locally rather than loading it from the
internet, so the whole system runs completely offline.

APScheduler runs our background jobs, and pytest runs our automated tests.
""",
    )

    # ── 9 ────────────────────────────────────────────────────────────
    s = _blank(prs)
    heading(s, "9", "Datasets")
    table(
        s,
        ["Dataset", "Rows", "Features", "What it predicts"],
        [
            ["Telco Customer Churn", "7,043", "19", "Will this customer cancel?"],
            ["Adult Census Income", "45,222", "13", "Does this person earn over $50K?"],
        ],
        col_widths=[3.0, 1.3, 1.5, 4.2],
        size=15,
        top=Inches(1.95),
    )
    cues(
        s,
        [
            "Split three ways: baseline, test, holdout",
            "Baseline is what all drift is compared against",
            "Holdout is kept aside for the simulator",
        ],
        top=Inches(3.5),
        size=20,
    )
    notes(
        s,
        """
We used two public datasets.

The first is Telco Customer Churn — about seven thousand telecom customers, with
nineteen pieces of information about each one, and whether they cancelled. That
is our main demonstration.

The second is Adult Census Income, about forty-five thousand records, predicting
whether a person earns more than fifty thousand dollars.

We cleaned both: removed rows with a missing answer, converted number columns
that were stored as text, and dropped ID columns, which carry no useful signal.

Then we split each dataset three ways. The baseline is the reference that all
future drift is measured against. The test set measures how good the model is.
And the holdout is kept completely aside, for the simulator to replay — so the
simulator uses real data the model has never seen.
""",
    )

    # ── 10 ───────────────────────────────────────────────────────────
    s = _blank(prs)
    heading(s, "10", "How It Works")
    table(
        s,
        ["Step", "What happens"],
        [
            ["1.  Data arrives", "A batch is uploaded, or the simulator sends one"],
            ["2.  Check quality", "Missing values, duplicates, outliers"],
            ["3.  Check drift", "Compare all 19 columns against the baseline"],
            ["4.  Check accuracy", "Score the batch, if the answers are known"],
            ["5.  Give a score", "Combine everything into one number out of 100"],
            ["6.  Raise alerts", "Warn, explain, and advise retraining"],
        ],
        col_widths=[2.6, 7.4],
        size=16,
        top=Inches(2.0),
    )
    notes(
        s,
        """
Every batch of data goes through six steps, always in this order.

Step one, the data arrives — either somebody uploads a file, or our simulator
sends one automatically.

Step two, we check the quality. Missing values, duplicate rows, outliers.

Step three, we check for drift, comparing all nineteen columns against the
baseline.

Step four, we score the batch with the model, but only if the true answers are
included.

Step five, we combine everything into a single health score out of 100.

Step six, we raise alerts, explain what changed, and advise retraining if
needed.

One thing worth mentioning — the order matters. We check quality before drift,
because a broken file would otherwise look exactly like a real change in
customers. And we check drift before accuracy, because drift can always be
measured, while accuracy often cannot.
""",
    )

    # ── 11 ───────────────────────────────────────────────────────────
    s = _blank(prs)
    heading(s, "11", "System Architecture")
    table(
        s,
        ["Layer", "What it does"],
        [
            ["Web pages", "What the user sees — charts, tables, alerts"],
            ["Django apps", "Handle requests and check permissions"],
            ["Services", "Coordinate one complete monitoring run"],
            ["Statistics engine", "Does the actual maths — no Django inside"],
            ["Scheduler", "Runs the simulator in the background"],
            ["Database", "Stores models, runs, alerts and settings"],
        ],
        col_widths=[3.0, 7.0],
        size=16,
        top=Inches(2.0),
    )
    notes(
        s,
        """
The system has six layers.

At the top, the web pages — what the user actually sees. Charts, tables, alerts.
Every chart can also be shown as a table, for accessibility.

Below that, the Django applications, which handle requests and check
permissions.

Then the service layer, which coordinates one complete monitoring run from start
to finish.

Then the most important part for us — the statistics engine. This is where the
actual mathematics lives, and we deliberately kept Django completely out of it.
That means we can test the statistics on their own, without a database and
without a web request. About two thousand lines of pure Python.

Below that, the scheduler that runs the simulator in the background, and finally
the database.
""",
    )

    # ── 12 ───────────────────────────────────────────────────────────
    s = _blank(prs)
    heading(s, "12", "Implementation")
    cues(
        s,
        [
            "Eight Django apps, each with one job",
            "Statistics engine kept separate from the website",
            "Permissions checked in three layers",
            "Five checks before a model file is accepted",
            "Simulator can create drift on demand",
            "344 automated tests, all passing",
        ],
        size=20,
    )
    notes(
        s,
        """
We built eight Django applications, each with one clear job. Accounts handles
login and roles. Registry handles models and their versions. Datasets handles
incoming data. Monitoring holds the statistics engine. Alerts handles warnings
and thresholds. Then dashboard, simulator, and a shared core.

The statistics engine is kept completely separate from the website, so it can be
tested on its own.

Permissions are checked in three layers. First, a filter decides which models
even exist for this user. Second, a role check. Third, a per-model permission
check.

Before we accept an uploaded model file, it must pass five checks — it must
load, it must have a predict method, it must successfully score fifty rows from
the baseline, it must return one answer per row, and every class it predicts
must be one the baseline actually contains.

And the simulator can create drift on demand, which is how we tested all of
this.

Everything is covered by 344 automated tests, and they all pass.
""",
    )

    # ── 13 ───────────────────────────────────────────────────────────
    s = _blank(prs)
    heading(s, "13", "Results")
    table(
        s,
        ["Version", "Accuracy", "Customers who left, correctly found"],
        [
            ["V1  Logistic Regression", "0.7878", "52.67%"],
            ["V2  Balanced Random Forest", "0.7523", "78.88%"],
            ["V3  Balanced Gradient Boosting", "0.7587", "69.25%"],
        ],
        col_widths=[3.6, 2.2, 4.2],
        size=16,
        top=Inches(2.0),
    )
    cues(
        s,
        [
            "V1 looks best — but misses half the leavers",
            "Drift found: 3 columns high, 3 moderate, 13 stable",
            "Contract went from 54.5% to 90.1% month-to-month",
            "Health score dropped to 70 out of 100",
        ],
        top=Inches(3.9),
        size=19,
    )
    notes(
        s,
        """
These are our results, measured on fourteen hundred customers the model had
never seen.

Look at the table. Version one has the highest accuracy — 0.7878. But look at
the last column. It finds only 52.67% of the customers who actually left.

It gets a high score by predicting "will not leave" for almost everybody. Since
most customers do stay, it is right most of the time — and completely useless,
because it fails at the one job it was built for.

Version three gives up about three points of accuracy and finds 69.25% instead.
That is nearly seventeen percentage points more real leavers. For a company
trying to keep customers, version three is clearly better.

This is exactly why our project tracks several measures instead of just
accuracy.

On drift, in our test run the system found three columns badly drifted, three
moderately, and thirteen stable. The Contract column had gone from 54.5%
month-to-month customers to 90.1% — a completely different kind of customer. The
health score fell to 70 out of 100, and the system raised a warning.
""",
    )

    # ── 14 ───────────────────────────────────────────────────────────
    s = _blank(prs)
    heading(s, "14", "Advantages")
    cues(
        s,
        [
            "Warns you before the true answers arrive",
            "Four measures, not one — stronger evidence",
            "Explains changes in plain English",
            "Missing answers handled correctly",
            "Old results never change",
            "Runs offline, with no licence cost",
            "Advises retraining, keeps a human in charge",
        ],
        size=20,
    )
    notes(
        s,
        """
The main advantages.

First and most important, it warns you early — from the incoming data alone,
before the true answers are available. That is weeks or months of warning.

Second, we use four measures instead of one. When several independent measures
agree, that is much stronger evidence than any single test.

Third, the system explains what changed in plain English, so the person who has
to act on it does not need to understand statistics.

Fourth, batches without answers are handled correctly as unknown, which removes
a whole category of false alarm.

Fifth, the settings are stored with each run, so old results can never be
quietly rewritten.

Sixth, it runs entirely on your own machine, offline, with no licence fee.

And finally, it advises retraining rather than doing it, which keeps a person
responsible for replacing a live model.
""",
    )

    # ── 15 ───────────────────────────────────────────────────────────
    s = _blank(prs)
    heading(s, "15", "Future Scope")
    cues(
        s,
        [
            "An API so systems can send data automatically",
            "Support for other kinds of prediction problems",
            "A bigger database for multiple users at once",
            "Email and message alerts",
            "Drift detection for images and text",
        ],
    )
    notes(
        s,
        """
There is more we would like to add.

Right now, data comes in by uploading a file or through the simulator. We would
like to add a proper API, so a company's own system can send data automatically
without anybody clicking anything.

At the moment we support yes-or-no predictions. We would like to extend that to
number predictions and to problems with more than two possible answers.

We would move from SQLite to PostgreSQL, so several people can use it heavily at
the same time.

We would send alerts by email and messaging, instead of only showing them on the
screen. The email code is written but not yet connected.

And finally, drift detection for images and text, which needs a different
approach from the one we use for tables.
""",
    )

    # ── 16 ───────────────────────────────────────────────────────────
    s = _blank(prs)
    heading(s, "16", "Conclusion")
    cues(
        s,
        [
            "Models fail silently — we made that visible",
            "Four measures, one clear health score",
            "Works even when answers are missing",
            "Roles enforced by the server, not the screen",
            "Tested on two datasets, 344 tests passing",
        ],
    )
    notes(
        s,
        """
To conclude.

We started with a problem that ordinary software monitoring does not cover — a
model that keeps running normally while quietly becoming wrong. We made that
failure visible.

DriftGuard checks drift with four measures, checks data quality, tracks
accuracy, and reports one clear health score with all its parts shown.

Two decisions set our work apart. We let the size of the change decide the drift
level, and use the statistical test only to confirm it. And we treat a batch
with no answers as unknown rather than zero, which prevents false alarms on
healthy models.

The three user roles are enforced by the server itself, not just hidden on the
screen — and we have tests that prove it by trying the forbidden actions
directly.

We validated everything on two public datasets and a drift simulator, with 344
automated tests passing.

Thank you. We are happy to take any questions.
""",
    )

    # ── 17 ───────────────────────────────────────────────────────────
    s = _blank(prs)
    heading(s, "17", "References")
    box = s.shapes.add_textbox(MARGIN, Inches(1.9), BODY_W, Inches(5.2))
    frame = box.text_frame
    frame.word_wrap = True
    refs = [
        "Sculley, D., Holt, G., Golovin, D., et al. (2015). Hidden Technical Debt in Machine Learning Systems. NeurIPS 28.",
        "Gama, J., Zliobaite, I., Bifet, A., Pechenizkiy, M., Bouchachia, A. (2014). A Survey on Concept Drift Adaptation. ACM Computing Surveys, 46(4).",
        "Lu, J., Liu, A., Dong, F., Gu, F., Gama, J., Zhang, G. (2019). Learning under Concept Drift: A Review. IEEE TKDE, 31(12).",
        "Massey, F. J. (1951). The Kolmogorov-Smirnov Test for Goodness of Fit. Journal of the American Statistical Association, 46(253).",
        "Lin, J. (1991). Divergence Measures Based on the Shannon Entropy. IEEE Transactions on Information Theory, 37(1).",
        "Siddiqi, N. (2006). Credit Risk Scorecards. John Wiley and Sons.",
        "Breck, E., Cai, S., Nielsen, E., Salib, M., Sculley, D. (2017). The ML Test Score. IEEE Big Data.",
        "Polyzotis, N., Roy, S., Whang, S. E., Zinkevich, M. (2017). Data Management Challenges in Production Machine Learning. ACM SIGMOD.",
        "Telco Customer Churn dataset. IBM Sample Data Sets, via Kaggle.",
        "Dua, D. and Graff, C. (2019). Adult Data Set. UCI Machine Learning Repository.",
    ]
    for i, ref in enumerate(refs):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        para.space_after = Pt(9)
        run = para.add_run()
        run.text = f"{i + 1}.   {ref}"
        run.font.size = Pt(13)
        run.font.name = "Calibri"
    notes(
        s,
        "These are our references. Verify each one against the source before submission.",
    )

    OUT.parent.mkdir(exist_ok=True)
    prs.save(OUT)
    return prs


if __name__ == "__main__":
    build()
    print(f"Wrote {OUT.name}")
