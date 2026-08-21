"""Generate the project presentation as a .pptx file.

Classic treatment: white ground, black text, one hairline rule under each
heading. No colour beyond that, because the deck is read on a projector and
printed in greyscale as often as it is presented.

Every figure on these slides comes from the built system — dataset row counts,
model metrics, drift measurements and test counts were all read out of the
project rather than written from memory. The one section that cannot be
verified from the code is the literature survey; see the note at the end of
this file.

    python scripts/build_ppt.py          # -> dist/DriftGuard_Presentation.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "dist" / "DriftGuard_Presentation.pptx"

BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x59, 0x59, 0x59)
RULE = RGBColor(0x00, 0x00, 0x00)

FONT = "Calibri"

# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGIN = Inches(0.85)
BODY_TOP = Inches(1.75)
BODY_W = SLIDE_W - (2 * MARGIN)


def _blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    return slide


def _text(
    slide, left, top, width, height, size, bold=False, color=BLACK, align=PP_ALIGN.LEFT
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    para = frame.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    return frame, run


def heading(slide, number, title):
    """Slide number, title, and the hairline rule beneath it."""
    _, run = _text(slide, MARGIN, Inches(0.55), BODY_W, Inches(0.4), 13, color=GREY)
    run.text = f"{number}"

    _, run = _text(slide, MARGIN, Inches(0.85), BODY_W, Inches(0.7), 30, bold=True)
    run.text = title

    line = slide.shapes.add_shape(1, MARGIN, Inches(1.58), BODY_W, Pt(1.2))
    line.fill.solid()
    line.fill.fore_color.rgb = RULE
    line.line.fill.background()
    line.shadow.inherit = False


def bullets(slide, items, top=BODY_TOP, size=17, gap=1.15):
    """items: list of str, or (str, [sub, sub]) for one level of nesting."""
    box = slide.shapes.add_textbox(MARGIN, top, BODY_W, SLIDE_H - top - Inches(0.6))
    frame = box.text_frame
    frame.word_wrap = True
    first = True

    for item in items:
        subs = []
        if isinstance(item, tuple):
            item, subs = item

        para = frame.paragraphs[0] if first else frame.add_paragraph()
        first = False
        para.space_after = Pt(7)
        para.line_spacing = gap
        run = para.add_run()
        run.text = f"•   {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = BLACK
        run.font.name = FONT

        for sub in subs:
            child = frame.add_paragraph()
            child.space_after = Pt(4)
            child.line_spacing = gap
            child.level = 1
            run = child.add_run()
            run.text = f"–   {sub}"
            run.font.size = Pt(size - 3)
            run.font.color.rgb = GREY
            run.font.name = FONT


def table(slide, headers, rows, top=BODY_TOP, col_widths=None, size=12.5):
    shape = slide.shapes.add_table(
        len(rows) + 1, len(headers), MARGIN, top, BODY_W, Inches(0.4)
    )
    tbl = shape.table

    if col_widths:
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(11.633 * w / total)

    def _style(cell, text, bold, fill):
        cell.text = str(text)
        para = cell.text_frame.paragraphs[0]
        run = para.runs[0] if para.runs else para.add_run()
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = WHITE if bold else BLACK
        run.font.name = FONT
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill

    for col, text in enumerate(headers):
        _style(tbl.cell(0, col), text, True, BLACK)

    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            _style(tbl.cell(r, c), text, False, WHITE)
    return tbl


def note(slide, text, top=None):
    top = top or (SLIDE_H - Inches(1.15))
    frame, run = _text(slide, MARGIN, top, BODY_W, Inches(0.7), 13, color=GREY)
    run.text = text
    frame.word_wrap = True


# ══════════════════════════════════════════════════════════════════════
# The deck
# ══════════════════════════════════════════════════════════════════════


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── Title ────────────────────────────────────────────────────────
    s = _blank(prs)
    _, run = _text(
        s,
        MARGIN,
        Inches(2.35),
        BODY_W,
        Inches(1.2),
        40,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    run.text = "DriftGuard"
    _, run = _text(
        s, MARGIN, Inches(3.25), BODY_W, Inches(1.0), 21, align=PP_ALIGN.CENTER
    )
    run.text = "A Multi-User Platform for Machine Learning\nModel Monitoring and Data Drift Detection"
    line = s.shapes.add_shape(1, Inches(5.4), Inches(4.45), Inches(2.5), Pt(1.2))
    line.fill.solid()
    line.fill.fore_color.rgb = RULE
    line.line.fill.background()
    line.shadow.inherit = False
    _, run = _text(
        s,
        MARGIN,
        Inches(4.85),
        BODY_W,
        Inches(1.2),
        15,
        color=GREY,
        align=PP_ALIGN.CENTER,
    )
    run.text = "Final Year Project\n\n[ Student names  ·  USN  ·  Guide  ·  Department  ·  College  ·  Year ]"

    # ── 1. Introduction ──────────────────────────────────────────────
    s = _blank(prs)
    heading(s, "1", "Introduction")
    bullets(
        s,
        [
            "A machine learning model is trained once on a fixed snapshot of data, then deployed into an environment that continues to change.",
            "The model does not adapt on its own. It keeps producing predictions with unchanged confidence even when the incoming data no longer resembles what it learned from.",
            (
                "Critically, this failure is silent:",
                [
                    "no exception is raised, no error is logged, no service goes down",
                    "the degradation is only discovered when the business outcome is already affected",
                ],
            ),
            "DriftGuard is a web-based platform that continuously monitors deployed models, detects statistical drift in incoming data, tracks predictive performance, and raises explained alerts.",
            "It is multi-user: three roles with distinct, server-enforced privileges reflect how this responsibility is divided in a real team.",
        ],
    )

    # ── 2. Motivation ────────────────────────────────────────────────
    s = _blank(prs)
    heading(s, "2", "Motivation")
    bullets(
        s,
        [
            "Conventional software fails loudly. A machine learning model fails quietly, which makes it far harder to detect.",
            (
                "Real consequences of undetected degradation:",
                [
                    "a churn model stops identifying customers who are about to leave, and retention budget is spent on the wrong people",
                    "a credit model built on one economic condition is applied to another, and risk is systematically misjudged",
                    "a fraud model misses new attack patterns it was never trained on",
                ],
            ),
            "Manual review does not scale. Checking distributions by hand across dozens of features, for several models, on every batch, is not sustainable.",
            "Accuracy alone is insufficient and often unavailable: ground truth for a prediction may only arrive weeks or months later.",
            "There is a need for a system that detects deterioration from the input data itself, before the true outcomes are known.",
        ],
    )

    # ── 3. Problem Statement ─────────────────────────────────────────
    s = _blank(prs)
    heading(s, "3", "Problem Statement")
    bullets(
        s,
        [
            "Machine learning models deployed in production degrade over time as the statistical properties of incoming data diverge from the training distribution, a phenomenon known as data drift.",
            "Existing deployments provide no mechanism to detect this degradation. Failures are silent and are typically identified only after measurable business loss.",
            (
                "The problem is compounded by three practical constraints:",
                [
                    "ground truth labels are delayed or unavailable, so accuracy cannot be relied upon as the primary signal",
                    "statistical significance tests alone are unreliable, because their outcome is dominated by sample size rather than by the practical importance of a change",
                    "monitoring in an organisation involves multiple people with different responsibilities and different levels of authority over a model",
                ],
            ),
            "The requirement is a multi-user monitoring platform that detects drift from input data alone, explains it in accessible terms, and enforces role-based access to model configuration.",
        ],
    )

    # ── 4. Objectives ────────────────────────────────────────────────
    s = _blank(prs)
    heading(s, "4", "Objectives")
    bullets(
        s,
        [
            "To detect data drift on every incoming batch using four complementary measures: Kolmogorov-Smirnov, Chi-Square, Population Stability Index and Jensen-Shannon Divergence.",
            "To assess incoming data quality: missing values, duplicate records, outliers, type mismatches and previously unseen categories.",
            "To track predictive performance over time and to handle correctly the common case in which ground truth is not yet available.",
            "To combine these signals into a single interpretable health score, with the contribution of each component shown.",
            "To generate deduplicated alerts and plain-language explanations of what changed in the data.",
            "To recommend, but never automatically perform, model retraining.",
            "To implement three roles with server-enforced authorisation and per-model access control.",
            "To provide a configurable data feed simulator that reproduces controlled drift for validation and demonstration.",
        ],
        size=15.5,
    )

    # ── 5. Literature Survey ─────────────────────────────────────────
    s = _blank(prs)
    heading(s, "5", "Literature Survey")
    table(
        s,
        ["Author(s) / Year", "Contribution", "Relevance to this work"],
        [
            [
                "Sculley et al., 2015",
                "Hidden technical debt in ML systems",
                "Identifies monitoring as a persistent, under-addressed cost of deployed models",
            ],
            [
                "Gama et al., 2014",
                "Survey on concept drift adaptation",
                "Establishes drift terminology and the detect-then-adapt framing",
            ],
            [
                "Lu et al., 2019",
                "Review of learning under concept drift",
                "Categorises drift detection methods; motivates using several measures together",
            ],
            [
                "Massey, 1951",
                "Kolmogorov-Smirnov goodness-of-fit test",
                "Basis of the significance test applied to numeric features",
            ],
            [
                "Lin, 1991",
                "Divergence measures based on Shannon entropy",
                "Defines Jensen-Shannon Divergence, used as the bounded magnitude measure",
            ],
            [
                "Siddiqi, 2006",
                "Credit risk scorecard methodology",
                "Source of the Population Stability Index and its conventional 0.10 / 0.25 bands",
            ],
            [
                "Breck et al., 2017",
                "The ML Test Score production-readiness rubric",
                "Motivates continuous validation of data and models after deployment",
            ],
            [
                "Polyzotis et al., 2017",
                "Data management challenges in production ML",
                "Establishes input data validation as a first-class production concern",
            ],
        ],
        col_widths=[2.0, 3.4, 4.6],
        size=11.5,
        top=Inches(1.85),
    )
    note(s, "Verify each citation against the source before submission.")

    # ── 6. Research Gap ──────────────────────────────────────────────
    s = _blank(prs)
    heading(s, "6", "Research Gap")
    bullets(
        s,
        [
            "Existing literature concentrates on drift detection algorithms in isolation, with limited attention to delivering them as a usable, multi-user system.",
            "Commercial monitoring platforms exist, but are proprietary, cloud-dependent and priced for enterprise deployment, placing them beyond academic and small-scale use.",
            (
                "Specific gaps addressed by this work:",
                [
                    "significance testing is commonly used as the sole criterion, although its outcome is dominated by sample size rather than by the magnitude of the change",
                    "batches without ground truth are frequently treated as scoring zero, which produces false alarms for models that are in fact healthy",
                    "explanations are generally presented as raw statistics rather than in language a non-specialist can act upon",
                    "role-based access control over monitoring configuration is rarely addressed, although in practice the person reading alerts is not the person who should be able to change thresholds",
                ],
            ),
        ],
    )

    # ── 7. Proposed Solution ─────────────────────────────────────────
    s = _blank(prs)
    heading(s, "7", "Proposed Solution")
    bullets(
        s,
        [
            "A self-hosted Django web application that monitors registered models against a fixed reference dataset, called the baseline.",
            (
                "Drift is classified by magnitude, and confirmed rather than decided by significance:",
                [
                    "PSI and Jensen-Shannon Divergence determine the drift band",
                    "K-S and Chi-Square tests confirm it; an unconfirmed result is reduced by one band rather than discarded",
                ],
            ),
            "All bin edges, quartiles and thresholds are derived from the baseline and never recomputed from the incoming batch, so a shifted batch cannot conceal its own shift.",
            "Batches without ground truth are recorded as unknown rather than zero, and the health score redistributes its weights across the components that remain measurable.",
            "Thresholds in force are stored with each run, so historical results remain exactly as they were originally assessed.",
            "Three roles enforced at the server: Administrator, Data Scientist and Analyst.",
        ],
    )

    # ── 8. Hardware / Software ───────────────────────────────────────
    s = _blank(prs)
    heading(s, "8", "Hardware and Software Requirements")
    table(
        s,
        ["Category", "Specification"],
        [
            ["Processor", "Intel Core i3 or equivalent, dual core (minimum)"],
            ["Memory", "4 GB RAM minimum, 8 GB recommended"],
            ["Storage", "2 GB free disk space"],
            ["Operating system", "Windows 10/11, macOS or Linux"],
            [
                "Language",
                "Python 3.11 (Django 5.0 does not support Python 3.12 or later)",
            ],
            ["Web framework", "Django 5.0.9"],
            ["Machine learning", "scikit-learn 1.5.2, joblib 1.4.2"],
            ["Data and statistics", "pandas 2.2.3, NumPy 1.26.4, SciPy 1.13.1"],
            ["Scheduling", "APScheduler 3.10.4 (in-process)"],
            ["Database", "SQLite 3 with write-ahead logging"],
            [
                "Front end",
                "Django templates, Chart.js 4.4.4, Alpine.js 3.14.1 (both bundled locally)",
            ],
            ["Static files", "WhiteNoise 6.7.0"],
            ["Testing", "pytest with pytest-django"],
        ],
        col_widths=[2.6, 7.4],
        size=12,
        top=Inches(1.8),
    )

    # ── 9. Datasets and Preprocessing ────────────────────────────────
    s = _blank(prs)
    heading(s, "9", "Datasets and Data Preprocessing")
    table(
        s,
        ["Dataset", "Total rows", "Features", "Baseline", "Test", "Holdout", "Target"],
        [
            [
                "Telco Customer Churn",
                "7,043",
                "19",
                "4,225",
                "1,409",
                "1,409",
                "Churn (Yes / No)",
            ],
            [
                "Adult Census Income",
                "45,222",
                "13",
                "27,133",
                "9,045",
                "9,044",
                "income (>50K / <=50K)",
            ],
        ],
        col_widths=[2.7, 1.3, 1.1, 1.2, 1.1, 1.2, 2.4],
        size=12,
        top=Inches(1.85),
    )
    bullets(
        s,
        [
            (
                "Preprocessing applied:",
                [
                    "removal of records with missing or malformed target values",
                    "conversion of numeric fields stored as text (for example TotalCharges)",
                    "exclusion of identifier columns, which carry no predictive signal",
                    "stratified three-way split preserving the class distribution across all partitions",
                ],
            ),
            "The holdout partition is reserved for the simulator, so replayed batches are data the model has never encountered.",
        ],
        top=Inches(3.15),
        size=15,
    )

    # ── 10. Methodology and Workflow ─────────────────────────────────
    s = _blank(prs)
    heading(s, "10", "Proposed Methodology and Workflow")
    table(
        s,
        ["Stage", "Operation", "Output"],
        [
            [
                "1. Ingestion",
                "Accept a batch by upload or from the simulator; validate the schema against the baseline",
                "Validated batch record",
            ],
            [
                "2. Data quality",
                "Missing values, duplicates, outliers by the baseline IQR rule, type mismatches, unseen categories",
                "Quality score out of 100",
            ],
            [
                "3. Drift detection",
                "Per feature: PSI, JSD, and K-S or Chi-Square; magnitude sets the band, the test confirms it",
                "Per-feature drift status",
            ],
            [
                "4. Performance",
                "Score the batch with the active model version, if ground truth is present",
                "Accuracy, precision, recall, F1",
            ],
            [
                "5. Health score",
                "Weighted composite of performance, drift, quality and stability",
                "Single score and band",
            ],
            [
                "6. Alerting",
                "Evaluate rules, deduplicate by model, category and feature; generate explanations",
                "Alerts and retraining advice",
            ],
        ],
        col_widths=[1.9, 5.6, 2.5],
        size=11.5,
        top=Inches(1.85),
    )
    note(
        s,
        "Quality precedes drift because a corrupted file would otherwise be indistinguishable from a genuine population change. "
        "Drift precedes performance because drift is measurable even when ground truth is absent.",
    )

    # ── 11. System Architecture ──────────────────────────────────────
    s = _blank(prs)
    heading(s, "11", "System Architecture")
    table(
        s,
        ["Layer", "Components", "Responsibility"],
        [
            [
                "Presentation",
                "Django templates, Chart.js, Alpine.js",
                "Role-aware interface; every chart also available as a table",
            ],
            [
                "Application",
                "accounts, registry, datasets, alerts, dashboard, simulator",
                "Views, authorisation decorators, request handling",
            ],
            [
                "Service",
                "monitoring.services, registry.services, alerts.services",
                "Orchestration; single convergence point for all ingestion",
            ],
            [
                "Analytical engine",
                "monitoring.engine (drift, quality, performance, health, explain)",
                "Pure Python; imports no Django, so it is independently testable",
            ],
            [
                "Scheduling",
                "APScheduler, in-process",
                "Drives simulator scenarios on a fixed interval",
            ],
            [
                "Persistence",
                "SQLite with WAL; file storage for artifacts and datasets",
                "Models, runs, alerts, immutable threshold snapshots",
            ],
        ],
        col_widths=[1.7, 4.0, 4.3],
        size=11.5,
        top=Inches(1.85),
    )
    note(
        s,
        "The analytical engine is deliberately isolated from the web framework: the statistical logic can be tested, "
        "reused and reasoned about without a database or an HTTP request.",
    )

    # ── 12. Implementation ───────────────────────────────────────────
    s = _blank(prs)
    heading(s, "12", "Implementation")
    bullets(
        s,
        [
            (
                "Eight functional Django applications:",
                [
                    "accounts (authentication, roles, per-model grants, login audit)",
                    "registry (models, versions, in-platform training, comparison)",
                    "datasets (baseline and batch ingestion, schema validation)",
                    "monitoring (the analytical engine and run orchestration)",
                    "alerts (rule evaluation, deduplication, thresholds, retraining advice)",
                    "dashboard, simulator, core (shared constants, validators, mixins)",
                ],
            ),
            "Authorisation enforced in three layers: a model visibility filter, a role decorator, and a per-model permission decorator.",
            "A model version must pass a five-check validation gate before it can be stored, and a version that fails validation cannot be activated.",
            "The simulator applies six configurable transformations to real held-out rows, expressed in baseline standard deviations so scenarios are scale-independent.",
        ],
        size=15.5,
    )
    note(
        s,
        "Approximately 15,600 lines of Python across 115 files, 32 templates, and 343 automated tests in 23 test modules. "
        "The analytical engine is about 2,000 lines and imports no Django.",
    )

    # ── 13. Expected Results ─────────────────────────────────────────
    s = _blank(prs)
    heading(s, "13", "Results")
    _, run = _text(s, MARGIN, Inches(1.78), BODY_W, Inches(0.35), 14, bold=True)
    run.text = "Model comparison on 1,409 held-out Telco records"
    table(
        s,
        ["Version", "Algorithm", "Accuracy", "Precision", "Recall", "F1"],
        [
            ["V1", "Logistic Regression", "0.7878", "0.6176", "0.5267", "0.5685"],
            ["V2", "Balanced Random Forest", "0.7523", "0.5221", "0.7888", "0.6283"],
            [
                "V3",
                "Balanced Gradient Boosting (active)",
                "0.7587",
                "0.5351",
                "0.6925",
                "0.6037",
            ],
        ],
        col_widths=[1.0, 4.0, 1.4, 1.4, 1.4, 1.4],
        size=12.5,
        top=Inches(2.2),
    )
    bullets(
        s,
        [
            "V1 attains the highest accuracy yet identifies only 52.67% of customers who actually churn; V3 identifies 69.25%, an improvement of 16.6 percentage points for a 2.9 point cost in accuracy.",
            "This demonstrates why accuracy alone is an inadequate criterion, and why the platform records several measures.",
            "Drift detection on a shifted batch: 3 features at high drift, 3 moderate, 13 stable, out of 19 monitored. Contract recorded PSI 0.721 against a 0.25 threshold, with the share of month-to-month contracts rising from 54.5% to 90.1%.",
            "Composite health score of 70/100, banded Warning, with drift contributing 37/100 and quality 92/100.",
            "343 automated tests pass, covering the statistical engine, the permission matrix, ingestion and error handling.",
        ],
        top=Inches(4.0),
        size=14,
    )

    # ── 14. Advantages ───────────────────────────────────────────────
    s = _blank(prs)
    heading(s, "14", "Advantages")
    bullets(
        s,
        [
            "Detects deterioration from input data alone, before ground truth becomes available, giving a genuine early warning.",
            "Four independent drift measures; agreement between them constitutes stronger evidence than any single test.",
            "Explanations are generated in plain language, so the person who must act on a finding need not interpret statistics.",
            "Unlabelled batches are handled correctly as unknown rather than as zero, eliminating an entire class of false alarm.",
            "Thresholds are recorded with each run, so historical results are immutable and auditable.",
            "Role-based access control is enforced at the server; a denied model is indistinguishable from one that does not exist, so probing reveals nothing.",
            "Entirely self-hosted with no external service dependency, and no per-model licensing cost.",
            "Retraining is advised rather than performed, keeping a human accountable for replacing a production model.",
        ],
        size=15.5,
    )

    # ── 15. Future Scope ─────────────────────────────────────────────
    s = _blank(prs)
    heading(s, "15", "Future Scope")
    bullets(
        s,
        [
            "A REST ingestion endpoint with API-key authentication, allowing production systems to submit batches automatically rather than through file upload.",
            "Support for regression and multi-class problems; the present implementation targets binary classification.",
            "Migration to PostgreSQL and an external task queue, permitting multiple concurrent workers.",
            "Delivery of alert notifications by email and webhook; the email module is written but is not yet invoked by the alerting flow.",
            "Automated retraining pipelines, retained under explicit human approval.",
            "Drift detection for image and text data, which requires embedding-based rather than tabular measures.",
            "Root-cause analysis linking a detected drift to the upstream system or process that produced it.",
        ],
    )

    # ── 16. Conclusion ───────────────────────────────────────────────
    s = _blank(prs)
    heading(s, "16", "Conclusion")
    bullets(
        s,
        [
            "DriftGuard addresses a failure mode that conventional software monitoring does not cover: a model that continues to operate normally while becoming progressively less correct.",
            "The platform detects drift using four complementary measures, assesses data quality, tracks performance, and reports a single interpretable health score with its components exposed.",
            (
                "Two design decisions distinguish it from a routine application of drift tests:",
                [
                    "magnitude determines the drift band while significance only confirms it, because a p-value reflects sample size rather than practical importance",
                    "a batch without ground truth is recorded as unknown rather than zero, which prevents systematic false alarms on healthy models",
                ],
            ),
            "Role-based access control is enforced by the server rather than concealed in the interface, and is verified by automated tests that attempt forbidden operations directly by URL.",
            "The system was validated on two public datasets and a configurable drift simulator, with 343 automated tests passing.",
        ],
        size=15.5,
    )

    # ── 17. References ───────────────────────────────────────────────
    s = _blank(prs)
    heading(s, "17", "References")
    bullets(
        s,
        [
            "Sculley, D., Holt, G., Golovin, D., et al. (2015). Hidden Technical Debt in Machine Learning Systems. Advances in Neural Information Processing Systems 28.",
            "Gama, J., Zliobaite, I., Bifet, A., Pechenizkiy, M., Bouchachia, A. (2014). A Survey on Concept Drift Adaptation. ACM Computing Surveys, 46(4).",
            "Lu, J., Liu, A., Dong, F., Gu, F., Gama, J., Zhang, G. (2019). Learning under Concept Drift: A Review. IEEE Transactions on Knowledge and Data Engineering, 31(12).",
            "Massey, F. J. (1951). The Kolmogorov-Smirnov Test for Goodness of Fit. Journal of the American Statistical Association, 46(253).",
            "Lin, J. (1991). Divergence Measures Based on the Shannon Entropy. IEEE Transactions on Information Theory, 37(1).",
            "Siddiqi, N. (2006). Credit Risk Scorecards: Developing and Implementing Intelligent Credit Scoring. John Wiley and Sons.",
            "Breck, E., Cai, S., Nielsen, E., Salib, M., Sculley, D. (2017). The ML Test Score: A Rubric for ML Production Readiness. IEEE International Conference on Big Data.",
            "Polyzotis, N., Roy, S., Whang, S. E., Zinkevich, M. (2017). Data Management Challenges in Production Machine Learning. ACM SIGMOD.",
            "Telco Customer Churn dataset. IBM Sample Data Sets, available via Kaggle.",
            "Dua, D. and Graff, C. (2019). Adult Data Set. UCI Machine Learning Repository, University of California, Irvine.",
        ],
        size=13,
        gap=1.0,
    )

    OUT.parent.mkdir(exist_ok=True)
    prs.save(OUT)
    return prs


if __name__ == "__main__":
    deck = build()
    print(f"Wrote {OUT.relative_to(ROOT)}")
    print(f"{len(deck.slides.__iter__.__self__._sldIdLst)} slides")
