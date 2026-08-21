"""Generate the styled version of the deck.

Same words and figures as the plain deck — both render from deck_content — but
with a visual identity taken from the product itself rather than from a generic
template. DriftGuard's interface is built on a deep indigo with a traffic-light
status scale; the deck uses that indigo as its primary and the amber it shows
for a Warning state as its single accent.

Restraint is the point. One accent colour, used for the rule under each heading
and the bullet markers, against a mostly white ground. A deck for a viva is read
on a projector that washes colour out, so the contrast has to survive that.

    python scripts/build_ppt_styled.py    # -> dist/DriftGuard_Presentation_Styled.pptx
"""

from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402
from pptx.util import Emu, Inches, Pt  # noqa: E402

import deck_content as C  # noqa: E402

OUT = (
    Path(__file__).resolve().parent.parent
    / "dist"
    / "DriftGuard_Presentation_Styled.pptx"
)

# ── Palette ───────────────────────────────────────────────────────────
INK = RGBColor(0x14, 0x14, 0x2B)  # near-black, blue bias
INDIGO = RGBColor(0x31, 0x2E, 0x81)  # primary
INDIGO_DK = RGBColor(0x1E, 0x1B, 0x54)  # title-slide ground
TINT = RGBColor(0xF2, 0xF3, 0xFA)  # panel / alternate row
AMBER = RGBColor(0xB4, 0x53, 0x09)  # the one accent
GREY = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

SPINE_W = Inches(0.30)  # indigo band down the left edge
MARGIN = Inches(1.05)
BODY_W = SLIDE_W - MARGIN - Inches(0.85)


def _rect(slide, left, top, width, height, colour):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _text(
    slide,
    left,
    top,
    width,
    height,
    size,
    text,
    bold=False,
    colour=INK,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    para = frame.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = colour
    run.font.name = FONT
    return frame


def _slide(prs, spine=True):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    if spine:
        _rect(s, Inches(0), Inches(0), SPINE_W, SLIDE_H, INDIGO)
    return s


def heading(slide, number, title):
    """Number badge, title, and the amber rule that carries the accent."""
    badge = _rect(slide, MARGIN, Inches(0.62), Inches(0.62), Inches(0.62), INDIGO)
    frame = badge.text_frame
    frame.word_wrap = False
    para = frame.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = number
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT

    _text(
        slide,
        MARGIN + Inches(0.9),
        Inches(0.66),
        BODY_W - Inches(0.9),
        Inches(0.7),
        31,
        title,
        bold=True,
    )

    _rect(slide, MARGIN, Inches(1.48), Inches(1.5), Pt(3), AMBER)
    _rect(
        slide,
        MARGIN + Inches(1.5),
        Inches(1.48) + Pt(1),
        BODY_W - Inches(1.5),
        Pt(1),
        RGBColor(0xDD, 0xDD, 0xE6),
    )


def cues(slide, items, top=Inches(2.05), size=22):
    """Bullets whose marker carries the accent colour, and whose text does not."""
    box = slide.shapes.add_textbox(MARGIN, top, BODY_W, SLIDE_H - top - Inches(0.5))
    frame = box.text_frame
    frame.word_wrap = True

    for i, item in enumerate(items):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        para.space_after = Pt(15)
        para.line_spacing = 1.1

        marker = para.add_run()
        marker.text = "▪   "
        marker.font.size = Pt(size)
        marker.font.color.rgb = AMBER
        marker.font.name = FONT

        run = para.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = INK
        run.font.name = FONT


def table(slide, spec, top=Inches(2.05)):
    headers, rows = spec["headers"], spec["rows"]
    widths, size = spec["widths"], spec.get("size", 14)

    shape = slide.shapes.add_table(
        len(rows) + 1, len(headers), MARGIN, top, BODY_W, Inches(0.4)
    )
    tbl = shape.table
    tbl.first_row = True
    tbl.horz_banding = False

    total = sum(widths)
    for i, w in enumerate(widths):
        tbl.columns[i].width = Inches(Emu(BODY_W).inches * w / total)

    def style(cell, text, *, header, fill):
        cell.text = str(text)
        para = cell.text_frame.paragraphs[0]
        run = para.runs[0] if para.runs else para.add_run()
        run.font.size = Pt(size)
        run.font.bold = header
        run.font.color.rgb = WHITE if header else INK
        run.font.name = FONT
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill

    for c, text in enumerate(headers):
        style(tbl.cell(0, c), text, header=True, fill=INDIGO)

    for r, row in enumerate(rows, start=1):
        # Banding by hand: PowerPoint's own banding follows the theme, which
        # this deck does not use.
        fill = TINT if r % 2 else WHITE
        for c, text in enumerate(row):
            style(tbl.cell(r, c), text, header=False, fill=fill)
    return tbl


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H

    # ── Title: full-bleed indigo ─────────────────────────────────────
    s = _slide(prs, spine=False)
    _rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, INDIGO_DK)
    _rect(s, Inches(0), Inches(0), SLIDE_W, Pt(7), AMBER)

    _text(
        s,
        Inches(1.0),
        Inches(2.25),
        SLIDE_W - Inches(2.0),
        Inches(1.3),
        46,
        C.TITLE["name"],
        bold=True,
        colour=WHITE,
        align=PP_ALIGN.CENTER,
    )
    _text(
        s,
        Inches(1.0),
        Inches(3.35),
        SLIDE_W - Inches(2.0),
        Inches(1.1),
        22,
        C.TITLE["subtitle"],
        colour=RGBColor(0xC7, 0xC9, 0xE8),
        align=PP_ALIGN.CENTER,
    )
    _rect(s, Inches(6.17), Inches(4.75), Inches(1.0), Pt(3), AMBER)
    _text(
        s,
        Inches(1.0),
        Inches(5.15),
        SLIDE_W - Inches(2.0),
        Inches(1.3),
        15,
        C.TITLE["footer"],
        colour=RGBColor(0x9A, 0x9D, 0xC4),
        align=PP_ALIGN.CENTER,
    )
    notes(s, C.TITLE["notes"])

    # ── Sections 1-16 ────────────────────────────────────────────────
    for spec in C.SLIDES:
        s = _slide(prs)
        heading(s, spec["number"], spec["title"])

        after_table = Inches(2.05)
        if "table" in spec:
            table(s, spec["table"])
            rows = len(spec["table"]["rows"]) + 1
            after_table = Inches(2.05 + rows * 0.42 + 0.35)
        if "cues" in spec:
            top = after_table if "table" in spec else Inches(2.05)
            size = 20 if len(spec["cues"]) > 5 or "table" in spec else 22
            cues(s, spec["cues"], top=top, size=size)

        _text(
            s,
            SLIDE_W - Inches(1.3),
            SLIDE_H - Inches(0.62),
            Inches(0.6),
            Inches(0.35),
            11,
            spec["number"],
            colour=GREY,
            align=PP_ALIGN.RIGHT,
        )
        notes(s, spec["notes"])

    # ── 17. References ───────────────────────────────────────────────
    s = _slide(prs)
    heading(s, "17", "References")
    box = s.shapes.add_textbox(MARGIN, Inches(1.95), BODY_W, Inches(5.1))
    frame = box.text_frame
    frame.word_wrap = True
    for i, ref in enumerate(C.REFERENCES):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        para.space_after = Pt(9)
        number = para.add_run()
        number.text = f"{i + 1}.   "
        number.font.size = Pt(13)
        number.font.bold = True
        number.font.color.rgb = INDIGO
        number.font.name = FONT
        run = para.add_run()
        run.text = ref
        run.font.size = Pt(13)
        run.font.color.rgb = INK
        run.font.name = FONT
    _text(
        s,
        SLIDE_W - Inches(1.3),
        SLIDE_H - Inches(0.62),
        Inches(0.6),
        Inches(0.35),
        11,
        "17",
        colour=GREY,
        align=PP_ALIGN.RIGHT,
    )
    notes(s, C.REFERENCES_NOTES)

    OUT.parent.mkdir(exist_ok=True)
    prs.save(OUT)
    return prs


if __name__ == "__main__":
    build()
    print(f"Wrote {OUT.name}")
