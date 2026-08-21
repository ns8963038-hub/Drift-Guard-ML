"""Render a .pptx to PDF without PowerPoint, Keynote or LibreOffice.

This machine has no application that opens .pptx — not even Quick Look — so the
deck could be built but never looked at. Rather than install an office suite,
this reads the shape geometry straight out of the file and lays each slide out
as an absolutely-positioned HTML page, then prints it with headless Chrome,
which Playwright already provides.

Because every position, size and font size is taken from the .pptx itself, the
PDF shows what the file actually contains rather than an approximation of it.

    python scripts/pptx_to_pdf.py dist/DriftGuard_Presentation_Simple.pptx
"""

from __future__ import annotations

import html
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from pptx.util import Emu

ROOT = Path(__file__).resolve().parent.parent


def _in(value) -> float:
    return Emu(value).inches if value is not None else 0.0


def _runs_html(paragraph) -> str:
    parts = []
    for run in paragraph.runs:
        size = run.font.size.pt if run.font.size else 18
        weight = "700" if run.font.bold else "400"
        colour = "#000"
        if run.font.color and run.font.color.type is not None:
            try:
                colour = f"#{run.font.color.rgb}"
            except (AttributeError, ValueError):
                pass
        text = html.escape(run.text).replace("\n", "<br>")
        parts.append(
            f'<span style="font-size:{size}pt;font-weight:{weight};color:{colour}">'
            f"{text}</span>"
        )
    return "".join(parts) or "&nbsp;"


def _solid_fill(shape) -> str | None:
    """The shape's solid fill as a CSS colour, or None if it has none.

    Needed because a styled deck is mostly rectangles — a full-bleed title
    ground, a spine, a number badge, an accent rule. Without this the converter
    drew none of them, so white text on an indigo ground came out invisible.
    """
    try:
        fill = shape.fill
        if fill.type != MSO_FILL.SOLID:
            return None
        return f"#{fill.fore_color.rgb}"
    except (AttributeError, TypeError, ValueError):
        return None


def slide_html(slide) -> str:
    out = []
    for shape in slide.shapes:
        left, top = _in(shape.left), _in(shape.top)
        width, height = _in(shape.width), _in(shape.height)

        if shape.has_table:
            cols = [_in(c.width) for c in shape.table.columns]
            rows_html = []
            for r, row in enumerate(shape.table.rows):
                cells = []
                for ci, cell in enumerate(row.cells):
                    para = cell.text_frame.paragraphs[0]
                    size, bold, colour = 12.0, False, "#000"
                    if para.runs:
                        run = para.runs[0]
                        size = run.font.size.pt if run.font.size else 12
                        bold = bool(run.font.bold)
                        if run.font.color and run.font.color.type is not None:
                            try:
                                colour = f"#{run.font.color.rgb}"
                            except (AttributeError, ValueError):
                                pass
                    bg = ""
                    try:
                        if cell.fill.type == MSO_FILL.SOLID:
                            bg = f"background:#{cell.fill.fore_color.rgb};"
                        elif cell.fill.type == MSO_FILL.BACKGROUND:
                            bg = "background:#fff;"
                    except (AttributeError, TypeError, ValueError):
                        pass
                    cells.append(
                        f'<td style="width:{cols[ci]}in;font-size:{size}pt;'
                        f'font-weight:{"700" if bold else "400"};color:{colour};{bg}">'
                        f"{html.escape(cell.text)}</td>"
                    )
                rows_html.append("<tr>" + "".join(cells) + "</tr>")
            out.append(
                f'<table class="tbl" style="left:{left}in;top:{top}in;'
                f'width:{sum(cols)}in">{"".join(rows_html)}</table>'
            )
            continue

        # A filled rectangle: the title ground, the spine, a rule, a badge.
        fill = _solid_fill(shape)
        if fill:
            out.append(
                f'<div class="fill" style="left:{left}in;top:{top}in;'
                f'width:{width}in;height:{height}in;background:{fill}"></div>'
            )

        if shape.has_text_frame and shape.text_frame.text.strip():
            lines = []
            for para in shape.text_frame.paragraphs:
                align = {2: "center", 3: "right", 4: "justify"}.get(
                    para.alignment.value if para.alignment else None, "left"
                )
                after = para.space_after.pt if para.space_after else 0
                indent = 0.42 if para.level else 0
                spacing = para.line_spacing or 1.15
                lines.append(
                    f'<div style="text-align:{align};margin-bottom:{after}pt;'
                    f'margin-left:{indent}in;line-height:{spacing}">'
                    f"{_runs_html(para)}</div>"
                )
            # Text sitting on a filled shape is centred in it, as PowerPoint does.
            extra = (
                f"height:{height}in;display:flex;flex-direction:column;justify-content:center;"
                if fill
                else ""
            )
            out.append(
                f'<div class="tb" style="left:{left}in;top:{top}in;'
                f'width:{width}in;{extra}">{"".join(lines)}</div>'
            )

    return f'<section class="slide">{"".join(out)}</section>'


def convert(src: Path, dest: Path) -> Path:
    from playwright.sync_api import sync_playwright

    prs = Presentation(str(src))
    w, h = _in(prs.slide_width), _in(prs.slide_height)

    body = "\n".join(slide_html(s) for s in prs.slides)
    page = f"""<!doctype html><meta charset="utf-8">
<style>
  @page {{ size: {w}in {h}in; margin: 0; }}
  * {{ box-sizing: border-box; }}
  body {{ margin: 0; font-family: Calibri, Carlito, "Helvetica Neue", Arial, sans-serif; }}
  .slide {{
    position: relative; width: {w}in; height: {h}in;
    background: #fff; overflow: hidden; page-break-after: always;
  }}
  .tb {{ position: absolute; }}
  .fill {{ position: absolute; }}
  .tbl {{ position: absolute; border-collapse: collapse; }}
  .tbl td {{ border: 0.5pt solid #dcdce4; padding: 5pt 7pt; vertical-align: top; }}
</style>
{body}"""

    tmp = dest.with_suffix(".html")
    tmp.write_text(page, encoding="utf-8")

    with sync_playwright() as pw:
        browser = pw.chromium.launch()
        pg = browser.new_context().new_page()
        pg.goto(tmp.resolve().as_uri())
        pg.wait_for_load_state("networkidle")
        pg.pdf(path=str(dest), width=f"{w}in", height=f"{h}in", print_background=True)
        browser.close()

    tmp.unlink()
    return dest


if __name__ == "__main__":
    source = (
        Path(sys.argv[1])
        if len(sys.argv) > 1
        else ROOT / "dist" / "DriftGuard_Presentation_Simple.pptx"
    )
    if not source.is_absolute():
        source = ROOT / source
    target = source.with_suffix(".pdf")
    convert(source, target)
    print(f"Wrote {target.name}  ({target.stat().st_size / 1024:.0f} KB)")
